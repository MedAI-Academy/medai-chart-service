"""
═══════════════════════════════════════════════════════════════
THEME PATCH v2 — Context-Aware Color Swapping
═══════════════════════════════════════════════════════════════

KEY INSIGHT: Templates have dark backgrounds with light text.
Light theme needs:
  - FILL colors: dark → light  (0F172A → FFFFFF)
  - TEXT colors: light → dark   (FFFFFF → 1E293B)

We CANNOT blindly swap both directions — that creates conflicts
where white fills become dark and dark text becomes white.

SOLUTION: Split the swap map into fill_swaps and text_swaps
based on color luminance. Dark colors only swap in fills,
light colors only swap in text.
═══════════════════════════════════════════════════════════════
"""

from pptx.dml.color import RGBColor
from lxml import etree
import logging

log = logging.getLogger(__name__)

NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def rgb_to_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def luminance(hex_str: str) -> int:
    """Simple perceived brightness: 0 (black) to 765 (white)."""
    h = hex_str.upper().lstrip('#')
    return int(h[0:2], 16) + int(h[2:4], 16) + int(h[4:6], 16)


def colors_match(c1: str, c2: str, tolerance: int = 10) -> bool:
    r1, g1, b1 = int(c1[0:2], 16), int(c1[2:4], 16), int(c1[4:6], 16)
    r2, g2, b2 = int(c2[0:2], 16), int(c2[2:4], 16), int(c2[4:6], 16)
    return abs(r1-r2) <= tolerance and abs(g1-g2) <= tolerance and abs(b1-b2) <= tolerance


def _find_swap(hex_val: str, swap_map: dict) -> str | None:
    hex_val = hex_val.upper()
    for src, tgt in swap_map.items():
        if hex_val == src.upper():
            return tgt.upper()
    for src, tgt in swap_map.items():
        if colors_match(hex_val, src.upper(), tolerance=10):
            return tgt.upper()
    return None


def apply_theme(slide, color_swap: dict):
    """
    Apply theme with CONTEXT-AWARE swapping.
    
    Splits color_swap into:
    - fill_swaps: dark→light (for shape/cell backgrounds)
    - text_swaps: light→dark (for font colors)
    - border_swaps: all (for lines/borders)
    """
    if not color_swap:
        return

    # Split by luminance: dark sources → fill swaps, light sources → text swaps
    DARK_THRESHOLD = 384   # R+G+B < 384 = "dark"
    
    fill_swaps = {}    # dark → light (backgrounds)
    text_swaps = {}    # light → dark (text)
    border_swaps = {}  # midtones (borders)
    
    for src, tgt in color_swap.items():
        src_lum = luminance(src)
        tgt_lum = luminance(tgt)
        
        if src_lum < DARK_THRESHOLD:
            # Dark source → this is a background swap (dark bg → light bg)
            fill_swaps[src] = tgt
        elif src_lum > 600:
            # Very light source → this is a text swap (white text → dark text)
            text_swaps[src] = tgt
        else:
            # Midtone → border/muted text swap (both directions)
            border_swaps[src] = tgt
    
    # Borders apply to both fills and text
    all_fill = {**fill_swaps, **border_swaps}
    all_text = {**text_swaps, **border_swaps}
    
    swapped = 0
    
    for shape in slide.shapes:
        swapped += _apply_fills(shape, all_fill)
        swapped += _apply_text(shape, all_text)
        swapped += _apply_lines(shape, color_swap)  # lines use full map
        
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    swapped += _apply_cell_fill(cell, all_fill)
                    if cell.text_frame:
                        swapped += _apply_text_frame(cell.text_frame, all_text)
        
        # Group shapes
        try:
            if hasattr(shape, 'shapes'):
                for child in shape.shapes:
                    swapped += _apply_fills(child, all_fill)
                    swapped += _apply_text(child, all_text)
                    swapped += _apply_lines(child, color_swap)
        except Exception:
            pass
    
    # Slide background
    try:
        bg = slide.background
        if bg.fill and bg.fill.type is not None:
            swapped += _swap_fill(bg.fill, all_fill)
    except Exception:
        pass
    
    # Also patch defRPr (default text properties) via raw XML
    swapped += _patch_default_text_colors(slide, all_text)
    
    if swapped > 0:
        log.info(f"  Theme v2: {swapped} swaps (fills:{len(fill_swaps)} text:{len(text_swaps)} border:{len(border_swaps)})")


def _apply_fills(shape, swap_map: dict) -> int:
    """Swap shape fill colors (solid + gradient)."""
    swapped = 0
    try:
        if hasattr(shape, 'fill') and shape.fill.type is not None:
            swapped += _swap_fill(shape.fill, swap_map)
    except Exception:
        pass
    return swapped


def _swap_fill(fill, swap_map: dict) -> int:
    swapped = 0
    try:
        ft = fill.type
    except:
        return 0
    
    if ft == 1:  # SOLID
        try:
            rgb = fill.fore_color.rgb
            if rgb:
                h = rgb_to_hex(rgb)
                new = _find_swap(h, swap_map)
                if new:
                    fill.solid()
                    fill.fore_color.rgb = hex_to_rgb(new)
                    swapped += 1
        except:
            pass
    elif ft == 3:  # GRADIENT
        try:
            for stop in fill.gradient_stops:
                try:
                    rgb = stop.color.rgb
                    if rgb:
                        h = rgb_to_hex(rgb)
                        new = _find_swap(h, swap_map)
                        if new:
                            stop.color.rgb = hex_to_rgb(new)
                            swapped += 1
                except:
                    pass
        except:
            pass
    return swapped


def _apply_text(shape, swap_map: dict) -> int:
    """Swap text font colors."""
    swapped = 0
    try:
        if shape.has_text_frame:
            swapped += _apply_text_frame(shape.text_frame, swap_map)
    except:
        pass
    return swapped


def _apply_text_frame(tf, swap_map: dict) -> int:
    swapped = 0
    for para in tf.paragraphs:
        for run in para.runs:
            try:
                fc = run.font.color
                if fc and fc.rgb:
                    h = rgb_to_hex(fc.rgb)
                    new = _find_swap(h, swap_map)
                    if new:
                        fc.rgb = hex_to_rgb(new)
                        swapped += 1
            except:
                pass
    return swapped


def _apply_lines(shape, swap_map: dict) -> int:
    """Swap line/border colors (use full map)."""
    swapped = 0
    try:
        line = shape.line
        if line.fill and line.fill.type is not None:
            swapped += _swap_fill(line.fill, swap_map)
        elif line.color and line.color.rgb:
            h = rgb_to_hex(line.color.rgb)
            new = _find_swap(h, swap_map)
            if new:
                line.color.rgb = hex_to_rgb(new)
                swapped += 1
    except:
        pass
    return swapped


def _apply_cell_fill(cell, swap_map: dict) -> int:
    """Swap table cell fill via raw XML."""
    swapped = 0
    try:
        tcPr = cell._tc.find(f'{{{NS}}}tcPr')
        if tcPr is not None:
            solidFill = tcPr.find(f'{{{NS}}}solidFill')
            if solidFill is not None:
                srgb = solidFill.find(f'{{{NS}}}srgbClr')
                if srgb is not None:
                    h = srgb.get('val', '').upper()
                    new = _find_swap(h, swap_map)
                    if new:
                        srgb.set('val', new)
                        swapped += 1
    except:
        pass
    return swapped


def _patch_default_text_colors(slide, text_swap_map: dict) -> int:
    """
    Patch <a:defRPr> elements — default text formatting that python-pptx
    doesn't expose via the high-level API. These control the "empty textbox"
    default color and are a major source of dark-on-dark text.
    """
    swapped = 0
    try:
        slide_xml = slide._element
        for defRPr in slide_xml.iter(f'{{{NS}}}defRPr'):
            solidFill = defRPr.find(f'{{{NS}}}solidFill')
            if solidFill is not None:
                srgb = solidFill.find(f'{{{NS}}}srgbClr')
                if srgb is not None:
                    h = srgb.get('val', '').upper()
                    new = _find_swap(h, text_swap_map)
                    if new:
                        srgb.set('val', new)
                        swapped += 1
    except Exception as e:
        log.debug(f"defRPr patch error: {e}")
    return swapped
