"""
Native Shape Renderer — Adds data-driven shapes directly to PPTX slides.

Charts rendered as editable PowerPoint shapes (not PNGs):
  - Forest Plot: diamonds + CI lines (only when CI data available)
  - Waterfall: sorted vertical bars
  - Swimmer: sorted horizontal bars
  - ORR Bars: grouped bar comparison
  
Data sources (in order of confidence):
  1. User-provided values (manual input) → Tier 1
  2. Vision-extracted from publication PDF → Tier 2 (labeled "Reconstructed")
  3. Never AI-estimated → FORBIDDEN

All data-driven visualizations sorted by convention:
  - Waterfall: descending (best response → worst response)
  - Swimmer: descending (longest duration → shortest)
  - Forest: ordered by subgroup category (as in publication)
"""

import math
import re
import logging
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# ── medaccur colors ──
PURPLE = RGBColor(0x7C, 0x6F, 0xFF)
TEAL = RGBColor(0x22, 0xD3, 0xA5)
ROSE = RGBColor(0xFF, 0x5F, 0x7E)
GOLD = RGBColor(0xF5, 0xC8, 0x42)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
DARK = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _parse_num(val):
    if val is None or val == '': return None
    if isinstance(val, (int, float)): return float(val)
    try:
        return float(str(val).replace('%', '').replace(',', '.').strip())
    except:
        return None


def _parse_ci_text(ci_text):
    """Parse CI from text like '(0.36-0.93)' or '0.36–0.93' or '0.36 to 0.93'."""
    if not ci_text: return None, None
    s = str(ci_text).strip().strip('()')
    # Try different separators
    for sep in ['–', '-', '—', ' to ', ',']:
        parts = s.split(sep)
        if len(parts) == 2:
            lo = _parse_num(parts[0])
            hi = _parse_num(parts[1])
            if lo is not None and hi is not None:
                return lo, hi
    return None, None


# ═══════════════════════════════════════════════════════
# FOREST PLOT — Diamonds + CI Lines
# ═══════════════════════════════════════════════════════
def add_forest_plot_shapes(slide, content):
    """
    Add diamonds and CI lines to forest plot.
    
    CI lines ONLY drawn when ci_low and ci_high are available.
    If CI data missing → only diamond at HR point, no line.
    """
    # Chart area (log scale, matching template axis labels 0.2 – 2.0)
    chart_left = 5.4
    chart_right = 12.2
    chart_width = chart_right - chart_left
    log_min = math.log(0.2)
    log_max = math.log(2.0)
    log_range = log_max - log_min

    def hr_to_x(hr):
        if hr <= 0: hr = 0.01
        if hr > 10: hr = 10
        log_hr = math.log(hr)
        frac = max(0, min(1, (log_hr - log_min) / log_range))
        return chart_left + frac * chart_width

    # Row definitions: (prefix, y_position)
    rows = [
        ('sg_overall', 1.5),
        ('sg_age_lt65', 2.2), ('sg_age_gte65', 2.4),
        ('sg_prior_1', 3.1), ('sg_prior_2_3', 3.3), ('sg_prior_gte4', 3.5),
        ('sg_imid_refract', 4.2), ('sg_pi_refract', 4.4), ('sg_double_refract', 4.6),
        ('sg_ecog_0', 5.3), ('sg_ecog_1', 5.5),
        ('sg_iss_1', 6.2), ('sg_iss_2', 6.4), ('sg_iss_3', 6.6),
    ]

    shapes_added = 0
    for prefix, y_pos in rows:
        hr_raw = content.get(prefix + '_hr', '')
        # Strip "HR " / "Hazard Ratio " prefix: "HR 0.54 (0.15-1.95)" → "0.54"
        import re as _re
        hr_cleaned = _re.sub(r'^(HR|Hazard Ratio)\s*', '', str(hr_raw), flags=_re.IGNORECASE).strip()
        hr_match = _re.match(r'[\d.]+', hr_cleaned)
        hr = float(hr_match.group()) if hr_match else None
        # Auto-extract CI from parentheses if not separately provided
        ci_in_parens = _re.search(r'\(([^)]+)\)', str(hr_raw))
        if ci_in_parens and not content.get(prefix + '_ci'):
            content[prefix + '_ci'] = ci_in_parens.group(1)
        if hr is None or hr <= 0 or hr > 5:
            continue

        x_hr = hr_to_x(hr)
        y_center = y_pos + 0.09

        # ── CI line: ONLY if ci_low + ci_high available ──
        ci_low = _parse_num(content.get(prefix + '_ci_low'))
        ci_high = _parse_num(content.get(prefix + '_ci_high'))

        # Try parsing from combined CI text (e.g., "0.36-0.93")
        if ci_low is None or ci_high is None:
            ci_text = content.get(prefix + '_ci', '')
            ci_low_parsed, ci_high_parsed = _parse_ci_text(ci_text)
            if ci_low_parsed is not None:
                ci_low = ci_low_parsed
            if ci_high_parsed is not None:
                ci_high = ci_high_parsed

        if ci_low is not None and ci_high is not None and ci_low > 0 and ci_high > ci_low:
            x_low = hr_to_x(ci_low)
            x_high = hr_to_x(ci_high)
            line_width = x_high - x_low
            if line_width > 0.03:
                ci_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x_low), Inches(y_center - 0.01),
                    Inches(line_width), Inches(0.02)
                )
                ci_shape.fill.solid()
                ci_shape.fill.fore_color.rgb = DARK
                ci_shape.line.fill.background()
                # CI end caps
                for x_cap in [x_low, x_high]:
                    cap = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(x_cap - 0.005), Inches(y_center - 0.06),
                        Inches(0.01), Inches(0.12)
                    )
                    cap.fill.solid()
                    cap.fill.fore_color.rgb = DARK
                    cap.line.fill.background()

        # ── Diamond at HR point (always drawn) ──
        is_overall = 'overall' in prefix
        diamond_size = 0.18 if is_overall else 0.13
        color = PURPLE if is_overall else TEAL

        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(x_hr - diamond_size / 2), Inches(y_center - diamond_size / 2),
            Inches(diamond_size), Inches(diamond_size)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = color
        diamond.line.fill.background()
        shapes_added += 1

    # ── Red reference line at HR=1.0 — LAST shape = always on top ──
    x_ref = hr_to_x(1.0)
    # Full height from top of chart area to bottom
    ref_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_ref - 0.005), Inches(1.3),
        Inches(0.015), Inches(5.5)
    )
    ref_line.fill.solid()
    ref_line.fill.fore_color.rgb = RGBColor(0xEF, 0x44, 0x44)  # Red
    ref_line.line.fill.background()
    shapes_added += 1

    logger.info(f"Forest plot: {shapes_added} shapes (diamonds + CI + ref line)")
    return shapes_added


# ═══════════════════════════════════════════════════════
# WATERFALL PLOT — Sorted Vertical Bars
# ═══════════════════════════════════════════════════════
def add_waterfall_shapes(slide, content):
    """
    Add waterfall bars sorted descending (best response first → worst).
    
    Data source: content['patients'] = [{change_pct, response}, ...]
    Source can be: user input OR vision-extracted from publication.
    """
    patients = content.get('patients', [])
    if not patients:
        return 0

    # SORT: descending by change (most negative = best response first on right)
    # Convention: waterfall goes from worst (left) to best (right)
    patients_sorted = sorted(patients, key=lambda p: _parse_num(p.get('change', p.get('change_pct', 0))) or 0, reverse=True)

    # Chart area
    area_left = 0.8
    area_right = 12.5
    area_top = 1.8
    area_bottom = 6.0
    area_width = area_right - area_left
    area_height = area_bottom - area_top
    y_zero = area_top + area_height * 0.5  # 0% line at center

    n = len(patients_sorted)
    if n == 0: return 0
    bar_width = min(0.4, (area_width - 0.5) / n)
    gap = (area_width - n * bar_width) / (n + 1)

    max_change = max(abs(_parse_num(p.get('change', p.get('change_pct', 0))) or 0) for p in patients_sorted)
    if max_change < 10: max_change = 100
    scale = (area_height * 0.45) / max_change  # pixels per %

    shapes_added = 0
    for i, pt in enumerate(patients_sorted):
        val = _parse_num(pt.get('change', pt.get('change_pct', 0))) or 0
        resp = pt.get('response', pt.get('status', ''))

        # Color by response
        if resp in ('CR', 'Complete'):
            color = PURPLE
        elif resp in ('PR', 'Partial') or val <= -30:
            color = TEAL
        elif resp in ('SD', 'Stable') or (-30 < val <= 20):
            color = GOLD
        else:
            color = ROSE

        x = area_left + gap + i * (bar_width + gap)
        bar_h = abs(val) * scale

        if val < 0:
            y = y_zero
        else:
            y = y_zero - bar_h

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(bar_width), Inches(max(bar_h, 0.02))
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        shapes_added += 1

    # Reference lines (-30% PR, +20% PD)
    for ref_val, ref_color, ref_label in [(-30, TEAL, 'PR -30%'), (20, ROSE, 'PD +20%')]:
        ref_y = y_zero - ref_val * scale
        if area_top < ref_y < area_bottom:
            ref_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(area_left), Inches(ref_y),
                Inches(area_width), Inches(0.01)
            )
            ref_line.fill.solid()
            ref_line.fill.fore_color.rgb = ref_color
            ref_line.line.fill.background()

    # Zero line
    zero_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(area_left), Inches(y_zero),
        Inches(area_width), Inches(0.015)
    )
    zero_line.fill.solid()
    zero_line.fill.fore_color.rgb = DARK
    zero_line.line.fill.background()

    logger.info(f"Waterfall: {shapes_added} bars (sorted descending)")
    return shapes_added


# ═══════════════════════════════════════════════════════
# SWIMMER PLOT — Sorted Horizontal Bars
# ═══════════════════════════════════════════════════════
def add_swimmer_shapes(slide, content):
    """
    Add swimmer bars sorted descending (longest duration at top).
    
    Data source: content['patients'] = [{duration, response, ongoing}, ...]
    Source can be: user input OR vision-extracted from publication.
    """
    patients = content.get('patients', [])
    if not patients:
        return 0

    # SORT: descending by duration (longest at top)
    patients_sorted = sorted(patients, key=lambda p: _parse_num(p.get('duration', p.get('duration_months', p.get('months', 0)))) or 0, reverse=True)

    # Chart area
    area_left = 1.5
    area_right = 12.5
    area_top = 1.5
    area_bottom = 6.2
    area_width = area_right - area_left
    area_height = area_bottom - area_top

    n = len(patients_sorted)
    if n == 0: return 0
    bar_height = min(0.35, (area_height - 0.2) / n)
    gap = (area_height - n * bar_height) / (n + 1)

    max_dur = max(_parse_num(p.get('duration', p.get('duration_months', p.get('months', 1)))) or 1 for p in patients_sorted)
    x_scale = area_width / (max_dur * 1.1)  # 10% padding

    shapes_added = 0
    for i, pt in enumerate(patients_sorted):
        dur = _parse_num(pt.get('duration', pt.get('duration_months', pt.get('months', 0)))) or 0
        resp = pt.get('response', pt.get('status', ''))
        ongoing = pt.get('ongoing', False)

        # Color
        if resp in ('CR', 'Complete'):
            color = PURPLE
        elif resp in ('PR', 'Partial'):
            color = TEAL
        elif resp in ('SD', 'Stable'):
            color = SLATE
        else:
            color = ROSE

        y = area_top + gap + i * (bar_height + gap)
        bar_w = dur * x_scale

        # Bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(area_left), Inches(y),
            Inches(max(bar_w, 0.05)), Inches(bar_height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Ongoing arrow
        if ongoing and bar_w > 0.3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(area_left + bar_w - 0.05), Inches(y + bar_height * 0.15),
                Inches(0.25), Inches(bar_height * 0.7)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()

        # Patient label (left side)
        label = pt.get('id', pt.get('patient', f'Pt {i+1}'))
        tb = slide.shapes.add_textbox(
            Inches(0.2), Inches(y),
            Inches(1.2), Inches(bar_height)
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = str(label)[:15]
        p.font.size = Pt(7)
        p.font.color.rgb = SLATE
        p.alignment = PP_ALIGN.RIGHT

        shapes_added += 1

    logger.info(f"Swimmer: {shapes_added} bars (sorted by duration desc)")
    return shapes_added


# ═══════════════════════════════════════════════════════
# ORR BAR CHART — Drug vs Control
# ═══════════════════════════════════════════════════════
def add_orr_bars(slide, content):
    """Add grouped bar chart for ORR/response rate comparison."""
    drug_val = _parse_num(content.get('chart_bar_exp_value', content.get('orr_exp_value')))
    ctrl_val = _parse_num(content.get('chart_bar_ctrl_value', content.get('orr_ctrl_value')))
    drug_label = content.get('chart_bar_exp_label', content.get('experimental_arm', 'Drug'))
    ctrl_label = content.get('chart_bar_ctrl_label', content.get('control_arm', 'Control'))

    if drug_val is None:
        return 0

    # Chart area position (right column of Pivotal Study template)
    chart_left = 8.2
    chart_bottom = 5.8
    chart_height = 3.5
    max_val = 100
    bar_width = 1.6
    gap = 0.5

    shapes = 0

    # Drug bar
    drug_h = chart_height * (drug_val / max_val)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(chart_left + 0.3), Inches(chart_bottom - drug_h),
        Inches(bar_width), Inches(drug_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()

    # Drug value label
    tb = slide.shapes.add_textbox(
        Inches(chart_left + 0.3), Inches(chart_bottom - drug_h - 0.3),
        Inches(bar_width), Inches(0.25)
    )
    p = tb.text_frame.paragraphs[0]
    p.text = f'{drug_val:.0f}%'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    # Drug arm label
    tb2 = slide.shapes.add_textbox(
        Inches(chart_left + 0.3), Inches(chart_bottom + 0.05),
        Inches(bar_width), Inches(0.2)
    )
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = str(drug_label)[:20]
    p2.font.size = Pt(8)
    p2.font.color.rgb = SLATE
    p2.alignment = PP_ALIGN.CENTER
    shapes += 1

    # Control bar (if available)
    if ctrl_val is not None and ctrl_val > 0:
        ctrl_h = chart_height * (ctrl_val / max_val)
        x_ctrl = chart_left + 0.3 + bar_width + gap
        bar2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_ctrl), Inches(chart_bottom - ctrl_h),
            Inches(bar_width), Inches(ctrl_h)
        )
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = SLATE
        bar2.line.fill.background()

        tb3 = slide.shapes.add_textbox(
            Inches(x_ctrl), Inches(chart_bottom - ctrl_h - 0.3),
            Inches(bar_width), Inches(0.25)
        )
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = f'{ctrl_val:.0f}%'
        p3.font.size = Pt(12)
        p3.font.bold = True
        p3.font.color.rgb = SLATE
        p3.alignment = PP_ALIGN.CENTER

        tb4 = slide.shapes.add_textbox(
            Inches(x_ctrl), Inches(chart_bottom + 0.05),
            Inches(bar_width), Inches(0.2)
        )
        p4 = tb4.text_frame.paragraphs[0]
        p4.text = str(ctrl_label)[:20]
        p4.font.size = Pt(8)
        p4.font.color.rgb = SLATE
        p4.alignment = PP_ALIGN.CENTER
        shapes += 1

    return shapes


# ═══════════════════════════════════════════════════════
# DISPATCHER — Called from deck_renderer after slide copy
# ═══════════════════════════════════════════════════════
def add_chart_shapes(slide, layout, content):
    """
    Add native shapes to a slide based on layout type.
    Called AFTER the template is copied to the output deck.
    
    Returns number of shapes added, or 0 if no chart needed.
    """
    dispatch = {
        'FOREST_PLOT': add_forest_plot_shapes,
        'WATERFALL_PLOT': add_waterfall_shapes,
        'SWIMMER_PLOT': add_swimmer_shapes,
        'PIVOTAL_STUDIES': add_orr_bars,
    }

    renderer = dispatch.get(layout)
    if not renderer:
        return 0

    try:
        return renderer(slide, content)
    except Exception as e:
        logger.error(f"Shape rendering error for {layout}: {e}")
        return 0
