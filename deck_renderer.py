"""
Deck Renderer — python-pptx Template Clone+Swap Engine

Assembles a complete PPTX from:
  1. Recipe JSON (ordered list of {layout, content})
  2. Template PPTX files (25 medaccur templates with {{placeholders}})
  3. Chart PNGs rendered inline for data-driven slides
  4. Theme color swapping for Light/Dark/Teal/Slate variants

Flow:
  - For each slide in recipe:
    - Open the correct template PPTX
    - Replace {{placeholders}} in all text shapes
    - Apply theme color swap (if not default dark theme)
    - For chart slides: render chart PNG, embed at chart_area position
    - Append the processed slide to the output PPTX
"""

import io
import os
import copy
import json
import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from lxml import etree

from theme_patch import apply_theme

logger = logging.getLogger(__name__)

# ── Namespace map for OpenXML ──
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

# ── Template directory (set at init or via env) ──
TEMPLATE_DIR = os.environ.get('TEMPLATE_DIR', 'templates/medaccur')
MANIFEST = None


def load_manifest():
    """Load medaccur_manifest.json"""
    global MANIFEST
    manifest_path = os.path.join(TEMPLATE_DIR, 'medaccur_manifest.json')
    if os.path.exists(manifest_path):
        with open(manifest_path, 'r') as f:
            MANIFEST = json.load(f)
        logger.info(f"Loaded manifest v{MANIFEST.get('version', '?')} with {MANIFEST.get('total_layouts', '?')} layouts")
    else:
        logger.warning(f"Manifest not found at {manifest_path}")
        MANIFEST = {"layout_map": {}, "section_to_layout": {}}
    return MANIFEST


def get_template_path(layout_id):
    """Resolve layout ID → template PPTX file path."""
    if not MANIFEST:
        load_manifest()
    layout_map = MANIFEST.get('layout_map', {})
    entry = layout_map.get(layout_id, {})
    filename = entry.get('file', '')
    if not filename:
        logger.warning(f"No template file for layout '{layout_id}'")
        return None
    path = os.path.join(TEMPLATE_DIR, filename)
    if not os.path.exists(path):
        logger.warning(f"Template file not found: {path}")
        return None
    return path


def replace_placeholders(slide, content):
    """
    Replace {{placeholder}} text in all shapes on a slide.
    
    Strategy: join all runs in a paragraph, match placeholders,
    put result text in first run (preserving formatting), clear rest.
    Handles cases where PowerPoint splits {{placeholder}} across runs.
    """
    if not content:
        return

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            runs = paragraph.runs
            if not runs:
                continue

            # Join all runs to get full paragraph text
            full_text = ''.join(run.text for run in runs)
            
            # Check if any placeholder exists
            has_placeholder = '{{' in full_text and '}}' in full_text
            if not has_placeholder:
                continue

            # Replace all {{key}} with values from content
            new_text = full_text
            for key, value in content.items():
                placeholder = '{{' + key + '}}'
                if placeholder in new_text:
                    # Convert value to string, handle None
                    val_str = str(value) if value is not None else ''
                    new_text = new_text.replace(placeholder, val_str)

            # Only update if something changed
            if new_text != full_text:
                # Put entire new text in first run, clear others
                runs[0].text = new_text
                for run in runs[1:]:
                    run.text = ''


def replace_unfilled_placeholders(slide):
    """Remove any remaining {{...}} placeholders that weren't filled."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            runs = paragraph.runs
            if not runs:
                continue
            full_text = ''.join(run.text for run in runs)
            if '{{' in full_text and '}}' in full_text:
                import re
                cleaned = re.sub(r'\{\{[^}]+\}\}', '', full_text)
                runs[0].text = cleaned
                for run in runs[1:]:
                    run.text = ''


def enable_auto_shrink(slide):
    """
    Enable text auto-shrink (normAutofit) on ALL text shapes in a slide.
    This prevents text from overflowing shape boundaries by automatically
    reducing font size to fit. Critical for AI-generated text which is
    often longer than the template placeholders expected.
    """
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        body_pr = shape.text_frame._txBody.find(f'{{{ns}}}bodyPr')
        if body_pr is None:
            continue
        # Remove any existing autofit settings
        for child in list(body_pr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('normAutofit', 'spAutoFit', 'noAutofit'):
                body_pr.remove(child)
        # Add normAutofit (shrink text to fit)
        autofit = etree.SubElement(body_pr, f'{{{ns}}}normAutofit')
        autofit.set('fontScale', '50000')  # Allow shrinking down to 50%
        count += 1
    return count


def embed_png_in_slide(slide, prs, png_bytes, left, top, width, height):
    """Embed a PNG image into a slide at the given position (in inches)."""
    from pptx.util import Inches as In
    image_stream = io.BytesIO(png_bytes)
    slide.shapes.add_picture(
        image_stream,
        In(left), In(top), In(width), In(height)
    )


def copy_slide_xml(source_prs, source_slide, target_prs):
    """
    Copy a slide from source presentation to target presentation.
    
    Uses XML-level manipulation to preserve all formatting, shapes,
    backgrounds, and images from the template.
    """
    # Add a blank slide to target
    blank_layout = target_prs.slide_layouts[0]
    target_slide = target_prs.slides.add_slide(blank_layout)

    # ── Copy slide background ──
    src_bg = source_slide._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
    tgt_bg = target_slide._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
    if src_bg is not None:
        if tgt_bg is not None:
            target_slide._element.remove(tgt_bg)
        # Insert background as first child of cSld
        csld = target_slide._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
        if csld is not None:
            csld.insert(0, copy.deepcopy(src_bg))

    # ── Replace shape tree entirely ──
    src_sp_tree = source_slide.shapes._spTree
    tgt_sp_tree = target_slide.shapes._spTree

    # Remove all existing shapes from target
    children_to_remove = []
    for child in tgt_sp_tree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
            children_to_remove.append(child)
    for child in children_to_remove:
        tgt_sp_tree.remove(child)

    # Copy all shapes from source
    for child in src_sp_tree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
            tgt_sp_tree.append(copy.deepcopy(child))

    return target_slide


def render_km_from_content(content):
    """
    Render KM curve PNG from km_data in slide content.
    Uses charts/km_curve.py if available, returns PNG bytes or None.
    """
    km_data = content.get('km_data')
    if not km_data:
        return None
    if not km_data.get('arms'):
        return None

    try:
        from charts.km_curve import render_km_curve
        png_bytes = render_km_curve(km_data)
        logger.info(f"  KM curve rendered: {len(png_bytes)} bytes, {len(km_data['arms'])} arms")
        return png_bytes
    except ImportError:
        logger.warning("  charts/km_curve.py not available — KM curve skipped")
        return None
    except Exception as e:
        logger.warning(f"  KM curve render failed: {e}")
        return None


def render_deck(recipe, chart_renderer=None, shape_renderer=None):
    """
    Build a complete PPTX from recipe JSON + templates.
    
    Args:
        recipe: dict with 'metadata', 'slides' (list of {id, layout, content}),
                optional 'theme' with 'color_swap' dict
        chart_renderer: optional callable(layout, content) -> bytes (PNG) — for KM curves
        shape_renderer: optional callable(slide, layout, content) -> int — for native shapes
    
    Returns:
        BytesIO with the assembled PPTX
    """
    if not MANIFEST:
        load_manifest()

    metadata = recipe.get('metadata', {})
    slides = recipe.get('slides', [])

    # ── Theme: read color swap map ──
    theme = recipe.get('theme', {})
    color_swap = theme.get('color_swap', {})
    theme_id = theme.get('id', 'dark')
    if color_swap:
        logger.info(f"Theme: {theme.get('name', theme_id)} — {len(color_swap)} color swaps")
    else:
        logger.info(f"Theme: default (no color swaps)")

    if not slides:
        raise ValueError("Recipe has no slides")

    logger.info(f"Rendering deck: {len(slides)} slides for {metadata.get('drug', '?')}")

    # ── Create output presentation ──
    output = Presentation()
    output.slide_width = Inches(13.33)
    output.slide_height = Inches(7.5)

    # Track which layouts need native shapes vs PNG
    SHAPE_LAYOUTS = {
        'FOREST_PLOT', 'WATERFALL_PLOT', 'SWIMMER_PLOT', 'PIVOTAL_STUDIES',
    }
    PNG_LAYOUTS = {
        'KM_CURVE',  # KM uses matplotlib PNG (curves too complex for shapes)
    }

    # Chart area position for KM PNG embedding
    CHART_AREAS = {
        'KM_CURVE': {'left': 0.5, 'top': 1.8, 'width': 12.3, 'height': 4.5},
    }

    slides_added = 0

    for i, slide_spec in enumerate(slides):
        layout = slide_spec.get('layout', '')
        content = slide_spec.get('content', {})
        slide_id = slide_spec.get('id', f'slide_{i}')

        # Inject metadata into content for all slides
        for mk, mv in metadata.items():
            if mk not in content:
                content[mk] = mv

        # Find template
        template_path = get_template_path(layout)
        if not template_path:
            logger.warning(f"Skipping slide {slide_id}: no template for layout '{layout}'")
            continue

        try:
            # Open template PPTX
            template_prs = Presentation(template_path)
            if not template_prs.slides:
                logger.warning(f"Template {layout} has no slides")
                continue
            template_slide = template_prs.slides[0]

            # ── Step 1: Replace {{placeholders}} ──
            replace_placeholders(template_slide, content)
            replace_unfilled_placeholders(template_slide)

            # ── Step 1b: Auto-shrink text to prevent overflow ──
            shrunk = enable_auto_shrink(template_slide)

            # ── Step 2: Apply theme color swap ──
            if color_swap:
                apply_theme(template_slide, color_swap)

            # ── Step 3: Copy processed slide to output ──
            output_slide = copy_slide_xml(template_prs, template_slide, output)
            slides_added += 1

            # ── Step 4: Native shapes (Forest, Waterfall, Swimmer, ORR bars) ──
            if layout in SHAPE_LAYOUTS and shape_renderer:
                try:
                    added = shape_renderer(output_slide, layout, content)
                    if added:
                        logger.info(f"  {added} native shapes added for {layout}")
                except Exception as shape_err:
                    logger.warning(f"  Shape rendering failed for {layout}: {shape_err}")

            # ── Step 5: KM curve PNG ──
            if layout in PNG_LAYOUTS:
                km_png = None

                # Try 1: Render from km_data in content (auto-extracted or manual)
                km_png = render_km_from_content(content)

                # Try 2: Use chart_renderer callback (legacy)
                if not km_png and chart_renderer:
                    try:
                        km_png = chart_renderer(layout, content)
                    except Exception as chart_err:
                        logger.warning(f"  Chart renderer failed for {layout}: {chart_err}")

                # Embed PNG if we got one
                if km_png:
                    area = CHART_AREAS.get(layout, {'left': 0.5, 'top': 1.8, 'width': 12.3, 'height': 4.5})
                    embed_png_in_slide(
                        output_slide, output, km_png,
                        area['left'], area['top'], area['width'], area['height']
                    )
                    logger.info(f"  KM PNG embedded for {layout}")

            logger.info(f"  [{slides_added}/{len(slides)}] {slide_id} → {layout}")

        except Exception as e:
            logger.error(f"Error processing slide {slide_id} ({layout}): {e}")
            continue

    if slides_added == 0:
        raise ValueError("No slides were successfully rendered")

    logger.info(f"Deck complete: {slides_added} slides, theme: {theme_id}")

    # ── Save to BytesIO ──
    buf = io.BytesIO()
    output.save(buf)
    buf.seek(0)
    return buf
